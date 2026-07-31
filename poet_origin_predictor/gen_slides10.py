# -*- coding: utf-8 -*-
"""
Generate the 10-minute version of the Tang Poets' Origins talk as .pptx.

v2: He Zhizhang thread removed; adds a task/class-definition slide;
transformer stress-test kept. 12 slides, ~50s each. Visual language follows
Tang_Poets_Origins_slides: brick red / cream / ink serif headings. Charts are
re-drawn in the deck palette (transparent PNGs in slides10_assets/).
Condensed speaker script is embedded in each slide's notes.

Usage: python3 gen_slides10.py
Output: Tang_Poets_Origins_slides_10min_v2.pptx
"""
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "slides10_assets")
os.makedirs(ASSETS, exist_ok=True)

# ---- palette (same as make_figures SLIDE_*) --------------------------
RED = RGBColor(0xA6, 0x3C, 0x2A)
INK = RGBColor(0x2E, 0x27, 0x23)
DARK_BG = RGBColor(0x2B, 0x23, 0x20)
CREAM = RGBColor(0xF2, 0xEE, 0xE7)
CREAM_RED = RGBColor(0xF6, 0xEA, 0xE6)
GRAY = RGBColor(0x8C, 0x84, 0x78)
LIGHT = RGBColor(0xFA, 0xF8, 0xF4)
BLUE = RGBColor(0x2F, 0x5A, 0x8F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

MRED, MGRAY, MINK = "#A63C2A", "#8C8478", "#2E2723"
HEAD_FONT = "Georgia"
BODY_FONT = "Calibri"

plt.rcParams["font.sans-serif"] = ["WenQuanYi Zen Hei", "Noto Sans CJK JP"]
plt.rcParams["axes.unicode_minus"] = False


# ---- charts ----------------------------------------------------------
def _style(ax):
    ax.tick_params(colors=MINK, labelsize=13)
    for sp in ax.spines.values():
        sp.set_color(MGRAY)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)


def chart_models():
    names = ["Baseline", "LogReg", "SVM", "RF", "MLP"]
    accs = [0.53, 0.60, 0.64, 0.67, 0.69]
    colors = [MGRAY] + ["#5A4F47"] * 3 + [MRED]
    fig, ax = plt.subplots(figsize=(6.4, 4.2))
    ax.bar(names, accs, color=colors, width=0.62)
    ax.axhline(0.5, ls="--", c=MGRAY, lw=1.2)
    ax.set_ylim(0, 0.8)
    ax.set_ylabel("accuracy (South/North)", fontsize=13, color=MINK)
    for i, v in enumerate(accs):
        ax.text(i, v + 0.015, f"{v:.2f}", ha="center", fontsize=13, color=MINK)
    _style(ax)
    fig.tight_layout()
    fig.savefig(os.path.join(ASSETS, "models.png"), dpi=150, transparent=True)
    plt.close(fig)


def chart_grains():
    names = ["10 circuits", "3 macro-regions", "South/North"]
    f1 = [0.18, 0.43, 0.69]
    fig, ax = plt.subplots(figsize=(6.0, 4.2))
    ax.bar(names, f1, color=["#5A4F47", "#5A4F47", MRED], width=0.55)
    ax.set_ylim(0, 0.8)
    ax.set_ylabel("macro-F1", fontsize=13, color=MINK)
    for i, v in enumerate(f1):
        ax.text(i, v + 0.015, f"{v:.2f}", ha="center", fontsize=13, color=MINK)
    _style(ax)
    fig.tight_layout()
    fig.savefig(os.path.join(ASSETS, "grains.png"), dpi=150, transparent=True)
    plt.close(fig)


def chart_eras():
    names = ["初唐\nEarly", "盛唐\nHigh", "中唐\nMid", "晚唐\nLate"]
    accs = [0.41, 0.50, 0.53, 0.68]
    ns = [29, 28, 43, 25]
    fig, ax = plt.subplots(figsize=(6.4, 4.2))
    ax.plot(names, accs, "o-", ms=11, lw=2.4, color=MRED)
    ax.axhline(0.5, ls="--", c=MGRAY, lw=1.2)
    ax.text(3.02, 0.505, "chance", fontsize=11, color=MGRAY, va="bottom", ha="right")
    for i, (a, n) in enumerate(zip(accs, ns)):
        ax.text(i, a + 0.025, f"{a:.2f}\n(n={n})", ha="center", fontsize=12,
                color=MINK)
    ax.set_ylim(0.3, 0.8)
    ax.set_ylabel("South/North accuracy", fontsize=13, color=MINK)
    _style(ax)
    fig.tight_layout()
    fig.savefig(os.path.join(ASSETS, "eras.png"), dpi=150, transparent=True)
    plt.close(fig)


# ---- pptx helpers ----------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SLIDE_N = 12


def add_slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BG if dark else WHITE
    return s


def rect(s, x, y, w, h, color, radius=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp


def text(s, x, y, w, h, runs, size=14, color=INK, bold=False, font=BODY_FONT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15,
         space_after=6):
    """runs: str, or list of paragraphs; each paragraph is str or list of
    (txt, {bold, color, size, font, italic}) run tuples."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        if isinstance(para, str):
            para = [(para, {})]
        for txt, kw in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(kw.get("size", size))
            r.font.bold = kw.get("bold", bold)
            r.font.italic = kw.get("italic", False)
            r.font.name = kw.get("font", font)
            r.font.color.rgb = kw.get("color", color)
    return tb


def chip(s):
    c = rect(s, 12.35, 0.42, 0.55, 0.55, RED)
    tf = c.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "詩"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = WHITE


def header(s, eyebrow, title):
    text(s, 0.62, 0.42, 10.5, 0.35, eyebrow, size=12, color=RED, bold=True)
    text(s, 0.62, 0.72, 11.4, 0.9, title, size=30, color=INK, bold=True,
         font=HEAD_FONT)
    chip(s)


def footer(s, idx):
    text(s, 0.62, 7.08, 8.0, 0.3,
         "Chen & Liu · Regional Fingerprints in the Complete Tang Poems",
         size=9, color=GRAY)
    text(s, 11.8, 7.08, 0.95, 0.3, f"{idx} / {SLIDE_N}", size=9, color=GRAY,
         align=PP_ALIGN.RIGHT)


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


def card(s, x, y, w, h, color=CREAM):
    return rect(s, x, y, w, h, color, radius=True)


# ---- build charts ----------------------------------------------------
chart_models()
chart_grains()
chart_eras()

# =====================================================================
# 1. TITLE (dark)
# =====================================================================
s = add_slide(dark=True)
text(s, 0.62, 1.55, 9.5, 0.35, "ARXIV 2606.24093 · COMPUTATIONAL LITERARY HISTORY",
     size=12, color=RED, bold=True)
text(s, 0.62, 1.95, 9.8, 2.0, "Predicting Poets' Origins from Verse",
     size=44, color=CREAM, bold=True, font=HEAD_FONT)
text(s, 0.62, 3.95, 9.8, 0.5,
     [[("Regional linguistic fingerprints in the ", {}),
       ("Complete Tang Poems", {"italic": True}),
       ("  全唐詩", {})]],
     size=17, color=RGBColor(0xC9, 0xC2, 0xB8))
text(s, 0.62, 4.7, 8.0, 0.9,
     [[("Chi-Sheng Chen", {"bold": True}), ("   Harvard University", {"color": GRAY})],
      [("Hung-Yun Liu", {"bold": True}),
       ("   University of Washington", {"color": GRAY})]],
     size=14, color=CREAM)
text(s, 0.62, 6.55, 3.0, 0.4, "10-minute version", size=11, color=GRAY)
for i, ch in enumerate("文如其地"):
    text(s, 12.35, 1.15 + i * 1.05, 0.6, 0.9, ch, size=28,
         color=RGBColor(0xC9, 0xC2, 0xB8), align=PP_ALIGN.CENTER)
c = rect(s, 12.38, 5.6, 0.55, 0.55, RED)
p = c.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "詩"; r.font.size = Pt(20); r.font.bold = True
r.font.color.rgb = WHITE
notes(s, "Good morning. Today I'd like to ask a very old question with a very "
         "new instrument: can you tell where a poet came from, just by reading "
         "their poems? We take the Complete Tang Poems — 49,000 of them — and "
         "treat this as a prediction problem. (~30s)")

# =====================================================================
# 2. QUESTION + RQs
# =====================================================================
s = add_slide()
header(s, "THE QUESTION", "Can a poem betray where its poet came from?")
card(s, 0.62, 1.85, 12.1, 2.7)
text(s, 0.98, 2.15, 11.4, 2.2,
     [[("Literary historians have long debated regional schools 地域流派 in "
        "Tang poetry — ", {}),
       ("yet whether geographic origin shows up as quantifiable patterns in "
        "poetic language has resisted systematic study.", {"bold": True})],
      [("We treat it as a measurement problem: ", {}),
       ("read ~49,000 poems and predict each poet's home circuit 道 from "
        "text alone.", {"bold": True, "color": RED})]],
     size=15, color=INK, line_spacing=1.25, space_after=10)
rq = [("RQ1 · Prediction", "Can a model recover a poet's circuit of origin "
       "from their collected work alone?"),
      ("RQ2 · Signal", "Which features carry it — lexical choice, imagery, "
       "tonal patterns, themes?"),
      ("RQ3 · Persistence", "Does the fingerprint hold steady across the "
       "dynasty's poetic eras?")]
for i, (t_, b_) in enumerate(rq):
    x = 0.62 + i * 4.15
    card(s, x, 4.95, 3.95, 1.7)
    text(s, x + 0.25, 5.15, 3.5, 0.35, t_, size=12, color=RED, bold=True)
    text(s, x + 0.25, 5.52, 3.5, 1.05, b_, size=11, color=INK)
text(s, 0.62, 6.75, 12.1, 0.35,
     "Framed for interpretability throughout — features a literary "
     "historian can read and contest.", size=11, color=GRAY,
     align=PP_ALIGN.CENTER)
footer(s, 2)
notes(s, "The regional-schools debate is old; what's been missing is "
         "quantification. We treat it as a measurement problem: predict each "
         "poet's home circuit from text alone. Three RQs: prediction, signal, "
         "persistence. Design principle: interpretability — features a "
         "literary historian can read and contest. (~55s)")

# =====================================================================
# 3. DATA
# =====================================================================
s = add_slide()
header(s, "DATA", "49,000 poems, 357 poets, ten circuits 道")
stats = [("900", "volumes, aggregated end-to-end"),
         ("≈49,000", "poems parsed into poet-level corpora"),
         ("357", "poets located via CBDB · ≥5 poems each"),
         ("242", "poets in the South/North subset")]
for i, (num, lab) in enumerate(stats):
    x = 0.62 + i * 3.15
    text(s, x, 2.0, 3.0, 0.9, num, size=34,
         color=RED if i in (1, 2) else INK, bold=True, font=HEAD_FONT)
    text(s, x, 2.85, 2.9, 0.7, lab, size=11, color=GRAY)
card(s, 0.62, 3.9, 12.1, 1.15)
text(s, 0.95, 4.05, 11.5, 0.9,
     [[("Pipeline.  ", {"bold": True}),
       ("Texts from the Chinese Text Project 中國哲學書電子化計劃; birthplaces "
        "linked via the China Biographical Database (CBDB); poets mapped to "
        "the ten early-Tang administrative circuits 道.", {})]],
     size=12, color=INK)
card(s, 0.62, 5.3, 12.1, 1.15, color=CREAM_RED)
text(s, 0.95, 5.45, 11.5, 0.9,
     [[("Heavily imbalanced.  ", {"bold": True, "color": RED}),
       ("Jiangnan 江南 alone has 126 poets; Longyou 隴右 has 6 — so: balanced "
        "class weights throughout, macro-F1 reported alongside accuracy.", {})]],
     size=12, color=INK)
footer(s, 3)
notes(s, "Complete Tang Poems, 900 volumes, ~49k poems into poet-level "
         "corpora. Birthplaces via CBDB; filter to >=5 poems, unambiguous "
         "attribution: 357 poets across ten circuits, 242 South/North. Flag "
         "now: heavy imbalance — Jiangnan 126, Longyou 6 — hence balanced "
         "weights + macro-F1. (~50s)")

# =====================================================================
# 4. TASKS & CLASS DEFINITIONS (new)
# =====================================================================
s = add_slide()
header(s, "SETUP", "Three tasks, three grains of geography")
# --- South/North
card(s, 0.62, 1.8, 3.95, 4.35, color=CREAM_RED)
text(s, 0.9, 2.0, 3.4, 0.4, [[("South 南 vs North 北", {"bold": True,
     "size": 15, "font": HEAD_FONT})]], color=INK)
text(s, 0.9, 2.45, 3.4, 0.35, "the headline task · 242 poets", size=10.5,
     color=RED, bold=True)
text(s, 0.9, 2.85, 3.45, 3.1,
     [[("南  ", {"bold": True, "color": RED}),
       ("江南 · 淮南 · 劍南 · 嶺南 · 山南", {})],
      [("北  ", {"bold": True, "color": BLUE}),
       ("河北 · 河南 · 河東 · 關內 · 隴右", {})],
      [("The axis the literary-historical debate cares about.",
        {"size": 10.5, "color": GRAY, "italic": True})]],
     size=11.5, color=INK, space_after=8)
# --- Macro
card(s, 4.72, 1.8, 3.95, 4.35)
text(s, 5.0, 2.0, 3.4, 0.4, [[("3 macro-regions", {"bold": True,
     "size": 15, "font": HEAD_FONT})]], color=INK)
text(s, 5.0, 2.45, 3.4, 0.35, "core vs periphery axis", size=10.5,
     color=GRAY, bold=True)
text(s, 5.0, 2.85, 3.45, 3.1,
     [[("東南 Southeast  ", {"bold": True, "color": RED}),
       ("江南 · 淮南", {})],
      [("中原 Central Plains  ", {"bold": True, "color": BLUE}),
       ("河北 · 河南 · 河東 · 關內", {})],
      [("邊陲 Frontier  ", {"bold": True, "color": GRAY}),
       ("劍南 · 嶺南 · 山南 · 隴右", {})]],
     size=11.5, color=INK, space_after=8)
# --- 10 circuits
card(s, 8.82, 1.8, 3.95, 4.35)
text(s, 9.1, 2.0, 3.4, 0.4, [[("10 circuits 道", {"bold": True,
     "size": 15, "font": HEAD_FONT})]], color=INK)
text(s, 9.1, 2.45, 3.4, 0.35, "the full administrative map · hardest",
     size=10.5, color=GRAY, bold=True)
text(s, 9.1, 2.85, 3.45, 3.1,
     [[("關內 · 河南 · 河北 · 河東 · 江南 · 淮南 · 山南 · 劍南 · 嶺南 · 隴右",
        {})],
      [("Many small classes — six poets in the smallest; classes with < 5 "
        "poets are dropped from CV.", {"size": 10.5, "color": GRAY,
                                       "italic": True})]],
     size=11.5, color=INK, space_after=8)
text(s, 0.62, 6.35, 12.1, 0.6,
     [[("Two different geographic hypotheses, not one nested hierarchy — ",
        {"color": GRAY}),
       ("山南·劍南·嶺南 are “south” on the S/N axis but “frontier” on the "
        "core–periphery axis.", {"bold": True, "color": INK})]],
     size=11, align=PP_ALIGN.CENTER)
footer(s, 4)
notes(s, "Three grains. South/North — the headline: five southern vs five "
         "northern circuits, 242 poets. Three macro-regions — a different "
         "axis: Southeast (Jiangnan+Huainan), Central Plains, Frontier. Ten "
         "circuits — the full map, hardest. Note: S/N and macro are two "
         "different hypotheses, not a nested hierarchy — Shannan, Jiannan, "
         "Lingnan are 'south' on one axis, 'frontier' on the other. (~55s)")

# =====================================================================
# 5. METHOD
# =====================================================================
s = add_slide()
header(s, "METHOD", "Readable features, one fair test")
card(s, 0.62, 1.85, 5.95, 4.5)
text(s, 0.95, 2.1, 5.3, 4.1,
     [[("Character n-gram TF-IDF.  ", {"bold": True}),
       ("空山新雨後 → 空·山·新·雨·後 + 空山·山新·新雨·雨後. 1–2 grams, "
        "≤8,000 features. No word segmentation, no modern assumptions — "
        "characters are the native unit.", {})],
      [("", {})],
      [("Plus features a historian can read:  ", {"bold": True}),
       ("imagery classes (山·水·草木·鳥獸·天體), seasonal markers, allusion "
        "density, lexical richness, tonal patterns 平仄 (approximated).", {})]],
     size=13, color=INK)
card(s, 6.8, 1.85, 5.95, 4.5)
text(s, 7.1, 2.1, 5.4, 0.4, "The fair test", size=15, color=INK, bold=True,
     font=HEAD_FONT)
tests = [("Five models", "LogReg · Linear SVM · Random Forest · MLP · "
          "GuwenBERT 古文BERT"),
         ("Stratified 5-fold CV", "every poet held out exactly once"),
         ("Balanced class weights", "small circuits not drowned out by Jiangnan"),
         ("Baseline", "“always guess the biggest region” — 0.53 on "
          "South/North")]
for i, (t_, b_) in enumerate(tests):
    text(s, 7.1, 2.65 + i * 0.88, 5.4, 0.8,
         [[(t_ + " — ", {"bold": True}), (b_, {})]], size=12, color=INK)
footer(s, 5)
notes(s, "Workhorse: character n-gram TF-IDF — no segmentation, no modern "
         "linguistic assumptions. On top, readable features: imagery, "
         "seasons, allusion, lexical richness, approximate tone. Five "
         "models incl. GuwenBERT, stratified 5-fold CV, balanced weights; "
         "baseline 0.53. (~50s)")

# =====================================================================
# 6. RQ1 RESULT
# =====================================================================
s = add_slide()
header(s, "RESULT · RQ1", "Geography is legible from verse alone")
s.shapes.add_picture(os.path.join(ASSETS, "models.png"),
                     Inches(0.62), Inches(1.95), Inches(7.3))
text(s, 0.9, 6.55, 7.0, 0.35,
     "South vs North · stratified 5-fold CV · 242 poets", size=10,
     color=GRAY, align=PP_ALIGN.CENTER)
card(s, 8.3, 1.95, 4.4, 4.6, color=CREAM_RED)
text(s, 8.65, 2.3, 3.8, 1.0, "0.69", size=48, color=RED, bold=True,
     font=HEAD_FONT)
text(s, 8.65, 3.45, 3.75, 2.9,
     [[("accuracy on South vs North — against a 0.53 most-frequent "
        "baseline", {})],
      [("Macro-F1 0.69 vs 0.35", {"bold": True}),
       (" — not an artefact of class imbalance.", {})],
      [("A poet's origin is readable, well above chance, from text alone.",
        {"italic": True, "color": RED})]],
     size=12, color=INK)
footer(s, 6)
notes(s, "Can it be done? Yes. Baseline 0.53 -> LogReg 0.60, SVM 0.64, RF "
         "0.67, MLP 0.69. Macro-F1 0.35 -> 0.69, so not an imbalance "
         "artefact. Origin is readable from text alone. (~45s)")

# =====================================================================
# 7. EVERY GRAIN
# =====================================================================
s = add_slide()
header(s, "RESULT", "The signal survives at every grain")
s.shapes.add_picture(os.path.join(ASSETS, "grains.png"),
                     Inches(0.62), Inches(2.0), Inches(6.8))
text(s, 0.9, 6.55, 6.4, 0.35,
     "macro-F1 by task granularity · each ≈2× its most-frequent baseline",
     size=10, color=GRAY, align=PP_ALIGN.CENTER)
card(s, 7.9, 2.0, 4.8, 4.4)
text(s, 8.2, 2.3, 4.3, 0.8, "Coarse maps are sharper — but even the hardest "
     "map is readable", size=14, color=INK, bold=True, font=HEAD_FONT)
rows = [("0.69", "South/North — the axis of the historical debate"),
        ("0.43", "3 macro-regions — regional structure persists"),
        ("0.18", "10 circuits — ≈2× baseline despite six-poet classes")]
for i, (n_, b_) in enumerate(rows):
    text(s, 8.2, 3.45 + i * 0.85, 4.3, 0.8,
         [[(n_ + "  ", {"bold": True, "color": RED, "size": 14}),
           (b_, {})]], size=12, color=INK)
footer(s, 7)
notes(s, "Harder maps: South/North 0.69, three macro-regions 0.43, full "
         "ten-circuit task 0.18 — each roughly TWICE its own baseline. "
         "Ten-way with six-poet classes is brutal; doubling baseline there "
         "is the surprise. (~40s)")

# =====================================================================
# 8. CENTRE BLURS / JIANGNAN
# =====================================================================
s = add_slide()
header(s, "RESULT", "The centre blurs; Jiangnan stands apart")
s.shapes.add_picture(os.path.join(HERE, "figures_en",
                                  "fig11b_periphery_regline.png"),
                     Inches(0.62), Inches(1.9), Inches(6.1))
text(s, 0.8, 6.6, 5.9, 0.4,
     "Per-circuit identifiability (recall) vs. distance from Chang'an — "
     "the apparent rise is carried by Jiangnan alone.", size=9, color=RED,
     align=PP_ALIGN.CENTER)
text(s, 7.1, 2.0, 5.6, 2.2,
     [[("Near the court, one idiom.  ", {"bold": True}),
       ("Circuits around Chang'an and Luoyang are heavily confused — the "
        "shared court language erases local difference.", {})],
      [("Far from it, a voice of one's own.  ", {"bold": True, "color": RED}),
       ("Jiangnan 江南 is by far the most separable region — recall 0.71 on "
        "the ten-way task.", {})]],
     size=13, color=INK)
card(s, 7.1, 4.75, 5.6, 1.85)
text(s, 7.4, 4.95, 5.0, 1.5,
     [[("Not a distance law.  ", {"bold": True}),
       ("All six circuits: non-significant (r = 0.45, p = 0.37); drop "
        "Jiangnan and the slope reverses (r = −0.89). Two facts, not one "
        "gradient.", {})]],
     size=12, color=INK)
footer(s, 8)
notes(s, "Where is the signal on the map? Two facts: capitals blur into one "
         "court idiom; Jiangnan stands apart at recall 0.71. Careful: NOT a "
         "distance law — full fit non-significant, slope reverses without "
         "Jiangnan. Two facts, not one gradient. (~55s)")

# =====================================================================
# 9. RQ2 IMAGERY
# =====================================================================
s = add_slide()
header(s, "RESULT · RQ2", "Imagery carries the region")
card(s, 0.62, 1.9, 5.95, 3.4, color=CREAM_RED)
text(s, 0.95, 2.15, 5.3, 0.5, [[("南  ", {"color": RED, "bold": True, "size": 18}),
     ("The South writes landscape", {"bold": True, "size": 16,
                                     "font": HEAD_FONT})]], color=INK)
text(s, 0.95, 2.8, 5.3, 2.3,
     [[("Mountain & water imagery 山水", {"bold": True}),
       (" — markedly more than in the north", {})],
      [("Buddhist & recluse vocabulary", {"bold": True}),
       (" — temples, retreat, withdrawal", {})],
      [("Landscape diction", {"bold": True}),
       (" — the discriminative characters of the south", {})]],
     size=12, color=INK)
card(s, 6.8, 1.9, 5.95, 3.4, color=RGBColor(0xE8, 0xEC, 0xF3))
text(s, 7.1, 2.15, 5.3, 0.5, [[("北  ", {"color": BLUE, "bold": True, "size": 18}),
     ("The North writes the court", {"bold": True, "size": 16,
                                     "font": HEAD_FONT})]], color=INK)
text(s, 7.1, 2.8, 5.3, 2.3,
     [[("Palace diction", {"bold": True}),
       (" — the built world of the capitals", {})],
      [("Gongti 宮體 motifs", {"bold": True}),
       (" — courtly “boudoir” poetry", {})],
      [("Court-idiom characters", {"bold": True}),
       (" — the discriminative characters of the north", {})]],
     size=12, color=INK)
card(s, 0.62, 5.55, 12.1, 1.1)
text(s, 0.95, 5.75, 11.5, 0.8,
     [[("Tonal patterns 平仄 carry comparatively weak signal — ", {}),
       ("the fingerprint lives in what poets see, not in how their verses "
        "scan.", {"bold": True})]],
     size=13, color=INK)
footer(s, 9)
notes(s, "RQ2: the South writes landscape — mountain/water imagery, "
         "Buddhist and recluse vocabulary. The North writes the court — "
         "palace diction, gongti motifs. Tone is weak. The fingerprint "
         "lives in what poets SEE, not how verses scan. (~45s)")

# =====================================================================
# 10. RQ3 ERAS
# =====================================================================
s = add_slide()
header(s, "RESULT · RQ3", "The fingerprint waxes and wanes with the dynasty")
s.shapes.add_picture(os.path.join(ASSETS, "eras.png"),
                     Inches(0.62), Inches(2.0), Inches(6.6))
labels = [("初唐 0.41", "southern poets assimilate — all 7 full-model errors "
           "run south → north"),
          ("盛唐 0.50", "at chance — maximum integration at the empire's "
           "height"),
          ("中唐 0.53", "after the An Lushan rebellion, regional voices begin "
           "to return"),
          ("晚唐 0.68", "as central authority weakens, separability peaks")]
for i, (t_, b_) in enumerate(labels):
    text(s, 7.6, 2.0 + i * 1.05, 5.1, 1.0,
         [[(t_ + "  ", {"bold": True, "color": RED}), (b_, {})]],
         size=12, color=INK)
text(s, 7.6, 6.25, 5.1, 0.6,
     "Regional distinctiveness waxes and wanes with the political "
     "integration of the dynasty.", size=12, color=INK, bold=True)
footer(s, 10)
notes(s, "RQ3: 0.41 early Tang — southerners wrote the northern court idiom, "
         "and all 7 full-model errors in that era run S->N. 0.50 at chance "
         "at the empire's height, 0.53 after An Lushan, 0.68 late Tang. "
         "Distinctiveness moves with political integration — now with "
         "numbers. (~55s)")

# =====================================================================
# 11. TRANSFORMERS + CAVEATS
# =====================================================================
s = add_slide()
header(s, "STRESS TEST & CAVEATS", "Do transformers read anything more?")
cols = [("0.62", "Naive fine-tuning", "GuwenBERT per fragment — below the "
         "n-gram models", INK),
        ("0.674", "Hierarchical frozen encoder", "pool fragments over a "
         "poet's corpus — ties the best classical models", BLUE),
        ("+ 0", "Hybrid BERT + TF-IDF", "combining the two adds nothing",
         INK)]
for i, (n_, t_, b_, c_) in enumerate(cols):
    x = 0.62 + i * 4.15
    card(s, x, 1.9, 3.95, 2.5)
    text(s, x + 0.3, 2.1, 3.4, 0.8, n_, size=30, color=c_, bold=True,
         font=HEAD_FONT)
    text(s, x + 0.3, 2.95, 3.4, 0.4, t_, size=13, color=INK, bold=True)
    text(s, x + 0.3, 3.4, 3.4, 0.9, b_, size=11, color=GRAY)
card(s, 0.62, 4.7, 12.1, 0.95, color=CREAM_RED)
text(s, 0.95, 4.87, 11.5, 0.7,
     [[("Character n-grams already capture the available regional signal — "
        "interpretability costs nothing.", {"bold": True, "color": RED})]],
     size=13, color=INK)
text(s, 0.62, 5.95, 12.1, 1.0,
     [[("Read with care:  ", {"bold": True}),
       ("small samples (357 / 242 / era subsets 25–43) · labels inherit "
        "CTP + CBDB noise · tone approximated · distance decay suggestive "
        "(r = 0.40, Mantel p ≈ 0.09). All four point the same way: more data "
        "would sharpen, not overturn.", {})]],
     size=12, color=INK)
footer(s, 11)
notes(s, "Wouldn't a transformer see more? Naive fine-tuning 0.62 — below "
         "n-grams. Hierarchical frozen encoder 0.674 — ties classical. "
         "Hybrid adds nothing. N-grams already capture the signal; "
         "interpretability costs nothing. Caveats in one breath: small "
         "samples, inherited label noise, approximate tone, suggestive "
         "geography — all point the same way. (~55s)")

# =====================================================================
# 12. CONCLUSION (dark)
# =====================================================================
s = add_slide(dark=True)
text(s, 0.62, 0.9, 9.0, 0.35, "CONCLUSION", size=12, color=RED, bold=True)
text(s, 0.62, 1.3, 10.6, 1.6,
     "A Tang poet's origin leaves a computationally detectable trace in "
     "their verse.", size=32, color=CREAM, bold=True, font=HEAD_FONT)
bullets = [("Carried chiefly by imagery and lexical choice",
            "mountains and waters in the south; palace and gongti diction "
            "in the north"),
           ("Not a simple distance law",
            "the capitals blur into one court idiom — Jiangnan alone keeps "
            "a voice of its own"),
           ("Modulated by history",
            "at chance amid High-Tang integration; strongest (0.68) as "
            "Late-Tang authority wanes")]
for i, (t_, b_) in enumerate(bullets):
    y = 3.25 + i * 0.85
    rect(s, 0.72, y + 0.12, 0.12, 0.12, RED)
    text(s, 1.05, y, 10.8, 0.8,
         [[(t_ + "  —  ", {"bold": True, "color": CREAM}),
           (b_, {"color": RGBColor(0xC9, 0xC2, 0xB8)})]], size=14)
text(s, 0.62, 6.0, 11.5, 0.7,
     [[("Next:  ", {"bold": True, "color": RED}),
       ("stylistic similarity networks — does textual proximity track "
        "geography, or transcend it through literary influence?",
        {"color": RGBColor(0xC9, 0xC2, 0xB8)})]], size=12)
text(s, 0.62, 6.7, 11.0, 0.6,
     [[("謝謝  Thank you", {"bold": True, "size": 16, "color": CREAM}),
       ("   ·   arXiv:2606.24093   ·   Chen & Liu", {"color": GRAY,
                                                     "size": 11})]])
for i, ch in enumerate("文如其地"):
    text(s, 12.35, 1.3 + i * 0.85, 0.6, 0.75, ch, size=26,
         color=RGBColor(0xC9, 0xC2, 0xB8), align=PP_ALIGN.CENTER)
c = rect(s, 12.38, 5.0, 0.55, 0.55, RED)
p = c.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "詩"; r.font.size = Pt(20); r.font.bold = True
r.font.color.rgb = WHITE
notes(s, "Conclusion: origin leaves a detectable trace — carried by imagery "
         "and lexical choice, not a simple distance law, modulated by "
         "history. Next: stylistic similarity networks. 文如其地 — the text "
         "is like its land. Thank you. (~40s)")

out = os.path.join(HERE, "Tang_Poets_Origins_slides_10min_v2.pptx")
prs.save(out)
print("saved", out)
